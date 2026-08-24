#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""OOXML block maps and marked temporary copies for Audion Docs AI."""

from __future__ import annotations

from dataclasses import dataclass
import hashlib
import os
from pathlib import Path
import shutil
from typing import Any, Dict, Iterable, List, Tuple


SUPPORTED_EXTS = {".docx", ".pptx"}
MARKER_PREFIX = "AB"


@dataclass(frozen=True)
class ProjectPaths:
    root: Path
    input_dir: Path
    output_dir: Path
    logs_dir: Path
    work_dir: Path
    cache_dir: Path


def default_project_paths(root: Path) -> ProjectPaths:
    source_override = str(os.environ.get("AUDION_WORKBENCH_SOURCE") or "").strip()
    target_override = str(os.environ.get("AUDION_WORKBENCH_TARGET") or "").strip()
    return ProjectPaths(
        root=root,
        input_dir=Path(source_override).expanduser().resolve(strict=False) if source_override else root / "input",
        output_dir=Path(target_override).expanduser().resolve(strict=False) if target_override else root / "output",
        logs_dir=root / "logs",
        work_dir=root / "work",
        cache_dir=root / "cache",
    )


def ensure_project_dirs(paths: ProjectPaths) -> None:
    for path in [
        paths.input_dir,
        paths.output_dir,
        paths.logs_dir,
        paths.work_dir,
        paths.cache_dir,
        paths.work_dir / "rendered_pdf",
        paths.work_dir / "marked_ooxml",
        paths.work_dir / "extracted_pdf_text",
    ]:
        if path == paths.input_dir and path.exists() and path.is_file():
            continue
        path.mkdir(parents=True, exist_ok=True)


def is_supported_document(path: Path) -> bool:
    return path.suffix.lower() in SUPPORTED_EXTS and not path.name.startswith("~$")


def iter_documents(input_dir: Path, *, recursive: bool = True) -> List[Path]:
    if input_dir.is_file():
        return [input_dir] if is_supported_document(input_dir) else []
    pattern = "**/*" if recursive else "*"
    return sorted(path for path in input_dir.glob(pattern) if path.is_file() and is_supported_document(path))


def source_root(input_dir: Path) -> Path:
    return input_dir.parent if input_dir.is_file() else input_dir


def source_relative(source: Path, input_dir: Path) -> Path:
    return source.relative_to(source_root(input_dir))


def rel_parent(source: Path, input_dir: Path) -> Path:
    rel = source_relative(source, input_dir)
    return rel.parent if str(rel.parent) != "." else Path()


def marker_for(block_id: str) -> str:
    digest = hashlib.blake2s(block_id.encode("utf-8"), digest_size=5).hexdigest()
    return f"{MARKER_PREFIX}{digest}"


def _safe_text(text: str) -> str:
    return " ".join((text or "").split())


def build_block_map(source: Path, input_dir: Path) -> Dict[str, Any]:
    ext = source.suffix.lower()
    if ext == ".docx":
        blocks = _build_docx_blocks(source)
    elif ext == ".pptx":
        blocks = _build_pptx_blocks(source)
    else:
        raise ValueError(f"Unsupported document type: {source}")

    return {
        "source_relative_path": source_relative(source, input_dir).as_posix(),
        "source_path": str(source),
        "document_type": ext.lstrip("."),
        "marker_prefix": MARKER_PREFIX,
        "blocks": blocks,
    }


def _build_docx_blocks(source: Path) -> List[Dict[str, Any]]:
    from docx import Document

    doc = Document(str(source))
    blocks: List[Dict[str, Any]] = []

    para_index = 0
    for paragraph in doc.paragraphs:
        text = _safe_text(paragraph.text)
        if not text:
            continue
        para_index += 1
        block_id = f"docx_p_{para_index:04d}"
        blocks.append(
            {
                "block_id": block_id,
                "marker": marker_for(block_id),
                "kind": "paragraph",
                "object_type": "paragraph",
                "text": text,
                "technical_location": {
                    "part": "word/document.xml",
                    "xpath_hint": f"/w:document/w:body/w:p[{para_index}]",
                    "paragraph_index": para_index,
                },
            }
        )

    for table_index, table in enumerate(doc.tables, start=1):
        for row_index, row in enumerate(table.rows, start=1):
            for cell_index, cell in enumerate(row.cells, start=1):
                for paragraph_index, paragraph in enumerate(cell.paragraphs, start=1):
                    text = _safe_text(paragraph.text)
                    if not text:
                        continue
                    block_id = (
                        f"docx_tbl_{table_index:04d}_"
                        f"r{row_index:03d}_c{cell_index:03d}_p{paragraph_index:03d}"
                    )
                    blocks.append(
                        {
                            "block_id": block_id,
                            "marker": marker_for(block_id),
                            "kind": "table_cell_paragraph",
                            "object_type": "table_cell",
                            "text": text,
                            "technical_location": {
                                "part": "word/document.xml",
                                "xpath_hint": (
                                    f"/w:document/w:body/w:tbl[{table_index}]"
                                    f"/w:tr[{row_index}]/w:tc[{cell_index}]/w:p[{paragraph_index}]"
                                ),
                                "table_index": table_index,
                                "row_index": row_index,
                                "cell_index": cell_index,
                                "paragraph_index": paragraph_index,
                            },
                        }
                    )

    return blocks


def _build_pptx_blocks(source: Path) -> List[Dict[str, Any]]:
    from pptx import Presentation

    prs = Presentation(str(source))
    blocks: List[Dict[str, Any]] = []

    for slide_index, slide in enumerate(prs.slides, start=1):
        shape_index = 0
        for shape in slide.shapes:
            shape_index += 1
            if getattr(shape, "has_table", False):
                table = shape.table
                for row_index, row in enumerate(table.rows, start=1):
                    for cell_index, cell in enumerate(row.cells, start=1):
                        for paragraph_index, paragraph in enumerate(cell.text_frame.paragraphs, start=1):
                            text = _safe_text(paragraph.text)
                            if not text:
                                continue
                            block_id = (
                                f"pptx_s{slide_index:03d}_tbl_{shape_index:04d}_"
                                f"r{row_index:03d}_c{cell_index:03d}_p{paragraph_index:03d}"
                            )
                            blocks.append(
                                {
                                    "block_id": block_id,
                                    "marker": marker_for(block_id),
                                    "kind": "table_cell_paragraph",
                                    "object_type": "table_cell",
                                    "text": text,
                                    "technical_location": {
                                        "part": f"ppt/slides/slide{slide_index}.xml",
                                        "slide_index": slide_index,
                                        "shape_index": shape_index,
                                        "row_index": row_index,
                                        "cell_index": cell_index,
                                        "paragraph_index": paragraph_index,
                                    },
                                }
                            )
                continue

            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs, start=1):
                text = _safe_text(paragraph.text)
                if not text:
                    continue
                block_id = f"pptx_s{slide_index:03d}_shape_{shape_index:04d}_p{paragraph_index:03d}"
                blocks.append(
                    {
                        "block_id": block_id,
                        "marker": marker_for(block_id),
                        "kind": "shape_paragraph",
                        "object_type": "paragraph",
                        "text": text,
                        "technical_location": {
                            "part": f"ppt/slides/slide{slide_index}.xml",
                            "slide_index": slide_index,
                            "shape_index": shape_index,
                            "paragraph_index": paragraph_index,
                        },
                    }
                )

    return blocks


def create_marked_copy(source: Path, marked_path: Path, block_map: Dict[str, Any]) -> Path:
    marked_path.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(source, marked_path)

    ext = source.suffix.lower()
    if ext == ".docx":
        _mark_docx(marked_path, block_map["blocks"])
    elif ext == ".pptx":
        _mark_pptx(marked_path, block_map["blocks"])
    else:
        raise ValueError(f"Unsupported document type: {source}")

    return marked_path


def _add_docx_marker_run(paragraph: Any, marker: str) -> None:
    from docx.shared import Pt, RGBColor

    run = paragraph.add_run(marker)
    run.font.size = Pt(2)
    run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    paragraph._p.remove(run._r)
    # WordprocessingML requires w:pPr to remain the first child of w:p.
    # Inserting a marker before it makes Word ignore paragraph formatting and
    # can drastically change pagination in the marked render copy.
    insert_at = 1 if paragraph._p.pPr is not None else 0
    paragraph._p.insert(insert_at, run._r)


def _mark_docx(path: Path, blocks: List[Dict[str, Any]]) -> None:
    from docx import Document

    doc = Document(str(path))
    by_id = {block["block_id"]: block for block in blocks}

    para_index = 0
    for paragraph in doc.paragraphs:
        if not _safe_text(paragraph.text):
            continue
        para_index += 1
        block = by_id.get(f"docx_p_{para_index:04d}")
        if block:
            _add_docx_marker_run(paragraph, block["marker"])

    for table_index, table in enumerate(doc.tables, start=1):
        for row_index, row in enumerate(table.rows, start=1):
            for cell_index, cell in enumerate(row.cells, start=1):
                for paragraph_index, paragraph in enumerate(cell.paragraphs, start=1):
                    if not _safe_text(paragraph.text):
                        continue
                    block_id = (
                        f"docx_tbl_{table_index:04d}_"
                        f"r{row_index:03d}_c{cell_index:03d}_p{paragraph_index:03d}"
                    )
                    block = by_id.get(block_id)
                    if block:
                        _add_docx_marker_run(paragraph, block["marker"])

    doc.save(str(path))


def _prefix_pptx_paragraph(paragraph: Any, marker: str) -> None:
    from pptx.dml.color import RGBColor
    from pptx.util import Pt

    if paragraph.runs:
        paragraph.runs[0].text = marker + " " + paragraph.runs[0].text
        run = paragraph.runs[0]
    else:
        run = paragraph.add_run()
        run.text = marker + " "
    run.font.size = Pt(1)
    run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)


def _mark_pptx(path: Path, blocks: List[Dict[str, Any]]) -> None:
    from pptx import Presentation

    prs = Presentation(str(path))
    by_id = {block["block_id"]: block for block in blocks}

    for slide_index, slide in enumerate(prs.slides, start=1):
        shape_index = 0
        for shape in slide.shapes:
            shape_index += 1
            if getattr(shape, "has_table", False):
                table = shape.table
                for row_index, row in enumerate(table.rows, start=1):
                    for cell_index, cell in enumerate(row.cells, start=1):
                        for paragraph_index, paragraph in enumerate(cell.text_frame.paragraphs, start=1):
                            if not _safe_text(paragraph.text):
                                continue
                            block_id = (
                                f"pptx_s{slide_index:03d}_tbl_{shape_index:04d}_"
                                f"r{row_index:03d}_c{cell_index:03d}_p{paragraph_index:03d}"
                            )
                            block = by_id.get(block_id)
                            if block:
                                _prefix_pptx_paragraph(paragraph, block["marker"])
                continue

            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs, start=1):
                if not _safe_text(paragraph.text):
                    continue
                block_id = f"pptx_s{slide_index:03d}_shape_{shape_index:04d}_p{paragraph_index:03d}"
                block = by_id.get(block_id)
                if block:
                    _prefix_pptx_paragraph(paragraph, block["marker"])

    prs.save(str(path))


def block_paragraph_refs_docx(path: Path) -> Dict[str, Any]:
    from docx import Document

    doc = Document(str(path))
    refs: Dict[str, Any] = {}
    para_index = 0
    for paragraph in doc.paragraphs:
        if not _safe_text(paragraph.text):
            continue
        para_index += 1
        refs[f"docx_p_{para_index:04d}"] = paragraph

    for table_index, table in enumerate(doc.tables, start=1):
        for row_index, row in enumerate(table.rows, start=1):
            for cell_index, cell in enumerate(row.cells, start=1):
                for paragraph_index, paragraph in enumerate(cell.paragraphs, start=1):
                    if not _safe_text(paragraph.text):
                        continue
                    block_id = (
                        f"docx_tbl_{table_index:04d}_"
                        f"r{row_index:03d}_c{cell_index:03d}_p{paragraph_index:03d}"
                    )
                    refs[block_id] = paragraph
    return {"document": doc, "refs": refs}


def block_paragraph_refs_pptx(path: Path) -> Dict[str, Any]:
    from pptx import Presentation

    prs = Presentation(str(path))
    refs: Dict[str, Any] = {}
    for slide_index, slide in enumerate(prs.slides, start=1):
        shape_index = 0
        for shape in slide.shapes:
            shape_index += 1
            if getattr(shape, "has_table", False):
                table = shape.table
                for row_index, row in enumerate(table.rows, start=1):
                    for cell_index, cell in enumerate(row.cells, start=1):
                        for paragraph_index, paragraph in enumerate(cell.text_frame.paragraphs, start=1):
                            if not _safe_text(paragraph.text):
                                continue
                            block_id = (
                                f"pptx_s{slide_index:03d}_tbl_{shape_index:04d}_"
                                f"r{row_index:03d}_c{cell_index:03d}_p{paragraph_index:03d}"
                            )
                            refs[block_id] = paragraph
                continue
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs, start=1):
                if not _safe_text(paragraph.text):
                    continue
                refs[f"pptx_s{slide_index:03d}_shape_{shape_index:04d}_p{paragraph_index:03d}"] = paragraph
    return {"presentation": prs, "refs": refs}
