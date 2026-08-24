#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Reporting core lifted verbatim from Audion Docs AI's pipeline.

Everything here is deterministic post-processing: it turns raw model findings into
the audit table, the report document and the anchored copy. No LLM calls live in
this module - in the skill flow the agent itself produces the findings.
"""

from __future__ import annotations

import hashlib
import json
import os
import re
import shutil
from pathlib import Path
from typing import Any, Dict, Iterable, List

from .document_model import block_paragraph_refs_docx, block_paragraph_refs_pptx
from .render_map import build_human_location, render_entry_by_block


REPORT_COLUMNS_EN = [
    "Issue ID",
    "Human Location",
    "Page",
    "Object Type",
    "Table",
    "Row",
    "Cell",
    "Paragraph",
    "Quote",
    "Problem",
    "Recommendation",
    "Fix Mode",
    "Confidence",
    "Old Text",
    "New Text",
    "Block ID",
    "Technical Location",
]

REPORT_COLUMNS_RU = [
    "ID ошибки",
    "Человеческая локация",
    "Страница",
    "Тип объекта",
    "Таблица",
    "Строка",
    "Ячейка",
    "Абзац",
    "Цитата",
    "Проблема",
    "Рекомендация",
    "Режим правки",
    "Уверенность",
    "Старый текст",
    "Новый текст",
    "ID блока",
    "Техническая локация",
]

ANCHOR_COLOR = "C05600"

ERROR_ANCHOR_RE = re.compile(r"\s*⟦[^⟧]{1,80}⟧")

OFFICIAL_TITLE_MARKERS_RE = re.compile(
    r"\b("
    r"постановлен\w*|распоряжен\w*|решени\w*|приказ\w*|закон\w*|"
    r"нормативн\w*|правов\w*|положени\w*|регламент\w*|правил\w*|"
    r"кодекс\w*|ред\.|№"
    r")\b",
    re.IGNORECASE,
)

OFFICIAL_NESTED_SOURCE_RE = re.compile(
    r"\(\s*вместе\s+с\s+«"
    r"(положением|правилами|порядком|перечнем|уставом|регламентом)\b"
    r".+«[^»]{2,160}»\s*\)",
    re.IGNORECASE | re.DOTALL,
)

REPORT_VALUE_MAP_RU = {
    "paragraph": "абзац",
    "table_cell": "ячейка таблицы",
    "slide_text": "текст слайда",
    "safe_replace": "безопасная замена",
    "requires_review": "требуется проверка",
    "none": "без правки",
    "high": "высокая",
    "medium": "средняя",
    "low": "низкая",
    "review": "проверка",
    "ok": "успешно",
}

def normalize_report_lang(value: str | None = None) -> str:
    raw = (value or os.environ.get("AUDION_REPORT_LANG") or "en").strip().lower()
    if raw in {"ru", "rus", "russian", "рус", "русский"}:
        return "ru"
    return "en"

def report_columns(lang: str) -> List[str]:
    return REPORT_COLUMNS_RU if normalize_report_lang(lang) == "ru" else REPORT_COLUMNS_EN

def report_sheet_title(lang: str) -> str:
    return "Аудит" if normalize_report_lang(lang) == "ru" else "Audit"

def report_docx_labels(lang: str) -> Dict[str, str | List[str]]:
    if normalize_report_lang(lang) == "ru":
        return {
            "status": "Статус",
            "issues": "Ошибок",
            "headers": ["ID ошибки", "Человеческая локация", "Проблема", "Рекомендация"],
        }
    return {
        "status": "Status",
        "issues": "Issues",
        "headers": ["Issue ID", "Human Location", "Problem", "Recommendation"],
    }

def localized_report_value(value: Any, report_lang: str) -> Any:
    if normalize_report_lang(report_lang) != "ru":
        return value
    text = str(value or "")
    return REPORT_VALUE_MAP_RU.get(text.strip().lower(), value)

def localized_human_text(value: Any, report_lang: str) -> str:
    text = str(value or "")
    if normalize_report_lang(report_lang) == "ru":
        text = re.sub(r"\bCHECK\s*:", "ПРОВЕРКА:", text, flags=re.IGNORECASE)
    return text

def formatted_technical_location(value: Any, report_lang: str) -> str:
    technical = value if isinstance(value, dict) else {}
    if normalize_report_lang(report_lang) != "ru":
        return json.dumps(technical, ensure_ascii=False)
    labels = {
        "part": "часть",
        "xpath_hint": "XPath",
        "paragraph_index": "абзац",
        "table_index": "таблица",
        "row_index": "строка",
        "cell_index": "ячейка",
        "slide_index": "слайд",
    }
    return "; ".join(
        f"{labels.get(key, key)}: {item}" for key, item in technical.items() if item not in (None, "")
    )

def sha256_file(path: Path) -> str:
    h = hashlib.sha256()
    with path.open("rb") as fh:
        for chunk in iter(lambda: fh.read(1024 * 1024), b""):
            h.update(chunk)
    return h.hexdigest()

def write_json(path: Path, payload: Dict[str, Any]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")

def read_json(path: Path) -> Dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8-sig", errors="replace"))

def normalize_issues(
    raw_issues: List[Dict[str, Any]],
    block_map: Dict[str, Any],
    render_map: Dict[str, Any],
) -> List[Dict[str, Any]]:
    blocks = {block["block_id"]: block for block in block_map.get("blocks", []) or []}
    renders = render_entry_by_block(render_map)
    normalized: List[Dict[str, Any]] = []

    for index, raw in enumerate(raw_issues, start=1):
        block_id = str(raw.get("block_id") or "").strip()
        if block_id not in blocks:
            continue
        block = blocks[block_id]
        if should_suppress_false_positive_issue(raw, block):
            continue
        render_entry = renders.get(block_id, {})
        technical = block.get("technical_location", {}) or {}
        issue_id = str(raw.get("issue_id") or f"E{index:03d}")
        normalized.append(
            {
                "issue_id": issue_id,
                "severity": raw.get("severity", "review"),
                "rule_id": raw.get("rule_id", ""),
                "human_location": build_human_location(block, render_entry),
                "page": render_entry.get("page"),
                "object_type": block.get("object_type") or block.get("kind", ""),
                "table": technical.get("table_index"),
                "row": technical.get("row_index"),
                "cell": technical.get("cell_index"),
                "paragraph": technical.get("paragraph_index"),
                "quote": raw.get("quote", block.get("text", "")),
                "problem": raw.get("problem") or raw.get("violation", ""),
                "recommendation": raw.get("recommendation") or raw.get("fix", ""),
                "fix_mode": raw.get("fix_mode", "requires_review"),
                "old_text": raw.get("old_text") or raw.get("source_text") or raw.get("replace_old", ""),
                "new_text": raw.get("new_text") or raw.get("replacement_text") or raw.get("replace_new", ""),
                "confidence": raw.get("confidence", ""),
                "status": raw.get("status", "requires_review"),
                "block_id": block_id,
                "technical_location": technical,
                "render_location": render_entry,
                "page_source": render_entry.get("page_source"),
                "page_confidence": render_entry.get("page_confidence"),
                "llm_raw_location": raw.get("location", ""),
            }
        )
    return normalized

def should_suppress_false_positive_issue(raw: Dict[str, Any], block: Dict[str, Any]) -> bool:
    return is_official_nested_quote_false_positive(raw, str(block.get("text") or ""))

def is_official_nested_quote_false_positive(raw: Dict[str, Any], block_text: str) -> bool:
    """Filter LLM-only quote complaints against exact official titles/citations."""
    problem = str(raw.get("problem") or raw.get("violation") or "")
    recommendation = str(raw.get("recommendation") or raw.get("fix") or "")
    combined = f"{problem} {recommendation}".lower()
    if "кавыч" not in combined:
        return False
    if not any(marker in combined for marker in ("сломан", "закрыва", "открыва", "незакрыт", "пар")):
        return False

    text = " ".join(str(block_text or "").split())
    if "«" not in text or "»" not in text:
        return False
    if not OFFICIAL_TITLE_MARKERS_RE.search(text):
        return False

    if text.count("«") == text.count("»"):
        return True

    return bool(OFFICIAL_NESTED_SOURCE_RE.search(text))

def build_empty_audit(source: Path, block_map: Dict[str, Any], render_map: Dict[str, Any], status: str) -> Dict[str, Any]:
    return {
        "source_relative_path": block_map.get("source_relative_path", source.name),
        "status": status,
        "issues": [],
        "meta": {
            "issue_source": "none",
            "note": "No LLM issue list was supplied to the pipeline; report artifacts are generated for smoke/structure validation.",
        },
    }

def block_text_for_llm(block: Dict[str, Any]) -> str:
    text = str(block.get("text") or "").strip()
    block_id = str(block.get("block_id") or block.get("id") or "").strip()
    kind = str(block.get("object_type") or block.get("kind") or "").strip()
    technical = block.get("technical_location", {}) or {}
    paragraph = technical.get("paragraph_index")
    table = technical.get("table_index")
    row = technical.get("row_index")
    cell = technical.get("cell_index")
    loc_bits = []
    if kind:
        loc_bits.append(f"type={kind}")
    if paragraph is not None:
        loc_bits.append(f"paragraph={paragraph}")
    if table is not None:
        loc_bits.append(f"table={table}")
    if row is not None:
        loc_bits.append(f"row={row}")
    if cell is not None:
        loc_bits.append(f"cell={cell}")
    suffix = f" ({', '.join(loc_bits)})" if loc_bits else ""
    return f"[BLOCK:{block_id}]{suffix}\n{text}"

def blocks_for_llm(block_map: Dict[str, Any]) -> List[Dict[str, Any]]:
    out: List[Dict[str, Any]] = []
    for block in block_map.get("blocks", []) or []:
        block_id = str(block.get("block_id") or "").strip()
        text = str(block.get("text") or "").strip()
        if not block_id or not text:
            continue
        out.append(
            {
                "id": block_id,
                "block_id": block_id,
                "kind": block.get("object_type") or block.get("kind", ""),
                "location": f"[BLOCK:{block_id}]",
                "text": block_text_for_llm(block),
            }
        )
    return out

def row_with_block_id(row: Dict[str, Any], block_map: Dict[str, Any]) -> Dict[str, Any]:
    import re

    out = dict(row)
    if out.get("block_id"):
        return out

    location = str(out.get("location") or "")
    match = re.search(r"\[BLOCK:([^\]]+)\]", location)
    if match:
        out["block_id"] = match.group(1).strip()
        return out

    quote = str(out.get("quote") or "").strip()
    if quote:
        for block in block_map.get("blocks", []) or []:
            text = str(block.get("text") or "")
            if quote in text:
                out["block_id"] = block.get("block_id")
                return out
    return out

def write_audit_table(out_path: Path, issues: List[Dict[str, Any]], report_lang: str = "en") -> None:
    from openpyxl import Workbook
    from openpyxl.styles import Alignment, Border, Font, PatternFill, Side

    out_path.parent.mkdir(parents=True, exist_ok=True)
    wb = Workbook()
    ws = wb.active
    ws.title = report_sheet_title(report_lang)
    ws.append(report_columns(report_lang))

    for issue in issues:
        technical = issue.get("technical_location", {})
        ws.append(
            [
                issue.get("issue_id", ""),
                issue.get("human_location", ""),
                issue.get("page", ""),
                localized_report_value(issue.get("object_type", ""), report_lang),
                issue.get("table", ""),
                issue.get("row", ""),
                issue.get("cell", ""),
                issue.get("paragraph", ""),
                issue.get("quote", ""),
                localized_human_text(issue.get("problem", ""), report_lang),
                localized_human_text(issue.get("recommendation", ""), report_lang),
                localized_report_value(issue.get("fix_mode", ""), report_lang),
                localized_report_value(issue.get("confidence", ""), report_lang),
                issue.get("old_text", ""),
                issue.get("new_text", ""),
                issue.get("block_id", ""),
                formatted_technical_location(technical, report_lang),
            ]
        )

    header_fill = PatternFill("solid", fgColor="1F4E78")
    header_font = Font(name="Tahoma", size=9, bold=True, color="FFFFFF")
    text_font = Font(name="Tahoma", size=9, color="1F1F1F")
    even_fill = PatternFill("solid", fgColor="F7F3EA")
    odd_fill = PatternFill("solid", fgColor="EFE7D8")
    thin = Side(style="thin", color="D9E2F3")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)

    for row in ws.iter_rows():
        for cell in row:
            cell.border = border
            cell.alignment = Alignment(horizontal="left", vertical="center", wrap_text=True, indent=1)
            if cell.row == 1:
                cell.font = header_font
                cell.fill = header_fill
                cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
            else:
                cell.font = text_font
                cell.fill = even_fill if cell.row % 2 == 0 else odd_fill

    for column_cells in ws.columns:
        max_len = max(len(str(cell.value or "")) for cell in column_cells)
        letter = column_cells[0].column_letter
        ws.column_dimensions[letter].width = min(max(max_len + 2, 8), 62)

    ws.freeze_panes = "A2"
    ws.auto_filter.ref = ws.dimensions
    ws.sheet_view.showGridLines = False

    for letter in ["I", "J", "K", "M"]:
        ws.column_dimensions[letter].width = max(ws.column_dimensions[letter].width or 0, 42)

    wb.save(str(out_path))

def write_audit_docx(out_path: Path, title: str, audit_payload: Dict[str, Any], report_lang: str = "en") -> None:
    from docx import Document
    from docx.shared import Pt

    out_path.parent.mkdir(parents=True, exist_ok=True)
    issues = audit_payload.get("issues", []) or []
    doc = Document()
    style = doc.styles["Normal"]
    style.font.name = "Tahoma"
    style.font.size = Pt(10)
    labels = report_docx_labels(report_lang)
    doc.add_heading(title, 0)
    doc.add_paragraph(
        f"{labels['status']}: {localized_report_value(audit_payload.get('status', ''), report_lang)}"
    )
    doc.add_paragraph(f"{labels['issues']}: {len(issues)}")

    table = doc.add_table(rows=1, cols=4)
    table.style = "Table Grid"
    headers = labels["headers"]
    for idx, header in enumerate(headers):
        table.cell(0, idx).text = header
    for issue in issues:
        row = table.add_row().cells
        row[0].text = str(issue.get("issue_id", ""))
        row[1].text = str(issue.get("human_location", ""))
        row[2].text = localized_human_text(issue.get("problem", ""), report_lang)
        row[3].text = localized_human_text(issue.get("recommendation", ""), report_lang)
    doc.save(str(out_path))

def write_annotated_document(
    source: Path,
    out_path: Path,
    block_map: Dict[str, Any],
    issues: List[Dict[str, Any]],
    log_path: Path,
) -> None:
    out_path.parent.mkdir(parents=True, exist_ok=True)
    issue_by_block: Dict[str, List[Dict[str, Any]]] = {}
    for issue in issues:
        issue_by_block.setdefault(issue.get("block_id", ""), []).append(issue)

    if source.suffix.lower() == ".docx":
        refs_payload = block_paragraph_refs_docx(source)
        doc = refs_payload["document"]
        refs = refs_payload["refs"]
        for block_id, block_issues in issue_by_block.items():
            paragraph = refs.get(block_id)
            if paragraph is None:
                continue
            for issue in block_issues:
                run = paragraph.add_run(f" ⟦{issue.get('issue_id', '')}⟧")
                run.bold = True
                try:
                    from docx.shared import RGBColor

                    run.font.color.rgb = RGBColor(0xC0, 0x56, 0x00)
                except Exception:
                    pass
        doc.save(str(out_path))
    elif source.suffix.lower() == ".pptx":
        refs_payload = block_paragraph_refs_pptx(source)
        prs = refs_payload["presentation"]
        refs = refs_payload["refs"]
        for block_id, block_issues in issue_by_block.items():
            paragraph = refs.get(block_id)
            if paragraph is None:
                continue
            for issue in block_issues:
                run = paragraph.add_run()
                run.text = f" ⟦{issue.get('issue_id', '')}⟧"
                try:
                    from pptx.dml.color import RGBColor

                    run.font.color.rgb = RGBColor(0xC0, 0x56, 0x00)
                    run.font.bold = True
                except Exception:
                    pass
        prs.save(str(out_path))
    else:
        shutil.copy2(source, out_path)

    write_json(
        log_path,
        {
            "source": str(source),
            "annotated": str(out_path),
            "issues": len(issues),
            "issue_ids": [issue.get("issue_id", "") for issue in issues],
            "status": "ok",
        },
    )

def _strip_anchor_runs(runs: Iterable[Any]) -> int:
    removed = 0
    for run in runs:
        text = str(getattr(run, "text", "") or "")
        matches = ERROR_ANCHOR_RE.findall(text)
        if not matches:
            continue
        run.text = ERROR_ANCHOR_RE.sub("", text)
        removed += len(matches)
    return removed

def _iter_docx_paragraphs(container: Any) -> Iterable[Any]:
    for paragraph in getattr(container, "paragraphs", []) or []:
        yield paragraph
    for table in getattr(container, "tables", []) or []:
        for row in table.rows:
            for cell in row.cells:
                yield from _iter_docx_paragraphs(cell)

def _iter_pptx_paragraphs(shapes: Any) -> Iterable[Any]:
    for shape in shapes:
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        yield paragraph
            continue
        if getattr(shape, "has_text_frame", False):
            for paragraph in shape.text_frame.paragraphs:
                yield paragraph
        nested_shapes = getattr(shape, "shapes", None)
        if nested_shapes is not None:
            yield from _iter_pptx_paragraphs(nested_shapes)

def strip_error_anchors_from_document(source: Path, out_path: Path) -> int:
    out_path.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(source, out_path)
    suffix = source.suffix.lower()
    removed = 0

    if suffix == ".docx":
        from docx import Document

        doc = Document(str(out_path))
        for paragraph in _iter_docx_paragraphs(doc):
            removed += _strip_anchor_runs(paragraph.runs)
        doc.save(str(out_path))
        return removed

    if suffix == ".pptx":
        from pptx import Presentation

        prs = Presentation(str(out_path))
        for slide in prs.slides:
            for paragraph in _iter_pptx_paragraphs(slide.shapes):
                removed += _strip_anchor_runs(paragraph.runs)
        prs.save(str(out_path))
        return removed

    return removed

def unanchored_output_path(source: Path) -> Path:
    stem = source.stem
    if stem.endswith("__annotated"):
        stem = stem[: -len("__annotated")]
    return source.with_name(f"{stem}__unanchored{source.suffix.lower()}")
