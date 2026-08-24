#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Recursive DOCX/PPTX audit pipeline with Office COM render-map support."""

from __future__ import annotations

import argparse
import hashlib
import json
import os
import re
import shutil
import sys
import threading
import time
from pathlib import Path
from typing import Any, Dict, Iterable, List

THIS_DIR = Path(__file__).resolve().parent
if str(THIS_DIR) not in sys.path:
    sys.path.insert(0, str(THIS_DIR))

from document_model import (
    block_paragraph_refs_docx,
    block_paragraph_refs_pptx,
    build_block_map,
    create_marked_copy,
    default_project_paths,
    ensure_project_dirs,
    iter_documents,
    rel_parent,
    source_relative,
)
from render.com_renderer import ConversionError, export_office_document_to_pdf, format_duration
from render.pdf_text import write_pdf_pages
from render.render_map import build_human_location, build_render_map, render_entry_by_block
from config_resolver import load_settings, resolve_api_key, resolve_model
from chunk_workers import DEFAULT_WORKERS, MAX_WORKERS, resolve_workers, run_jobs
from llm_audit_helpers import build_chunks, build_instructions, build_user_prompt, dedupe_rows, load_rules


REPORT_COLUMNS_EN = [
    "Issue ID",
    "Human Location",
    "Page",
    "Object Type",
    "Table",
    "Row",
    "Cell",
    "Paragraph",
    "Quote",
    "Problem",
    "Recommendation",
    "Fix Mode",
    "Confidence",
    "Old Text",
    "New Text",
    "Block ID",
    "Technical Location",
]

REPORT_COLUMNS_RU = [
    "ID ошибки",
    "Человеческая локация",
    "Страница",
    "Тип объекта",
    "Таблица",
    "Строка",
    "Ячейка",
    "Абзац",
    "Цитата",
    "Проблема",
    "Рекомендация",
    "Режим правки",
    "Уверенность",
    "Старый текст",
    "Новый текст",
    "ID блока",
    "Техническая локация",
]

ANCHOR_COLOR = "C05600"
ERROR_ANCHOR_RE = re.compile(r"\s*⟦[^⟧]{1,80}⟧")
OFFICIAL_TITLE_MARKERS_RE = re.compile(
    r"\b("
    r"постановлен\w*|распоряжен\w*|решени\w*|приказ\w*|закон\w*|"
    r"нормативн\w*|правов\w*|положени\w*|регламент\w*|правил\w*|"
    r"кодекс\w*|ред\.|№"
    r")\b",
    re.IGNORECASE,
)
OFFICIAL_NESTED_SOURCE_RE = re.compile(
    r"\(\s*вместе\s+с\s+«"
    r"(положением|правилами|порядком|перечнем|уставом|регламентом)\b"
    r".+«[^»]{2,160}»\s*\)",
    re.IGNORECASE | re.DOTALL,
)


def normalize_report_lang(value: str | None = None) -> str:
    raw = (value or os.environ.get("AUDION_REPORT_LANG") or "en").strip().lower()
    if raw in {"ru", "rus", "russian", "рус", "русский"}:
        return "ru"
    return "en"


def report_columns(lang: str) -> List[str]:
    return REPORT_COLUMNS_RU if normalize_report_lang(lang) == "ru" else REPORT_COLUMNS_EN


def report_sheet_title(lang: str) -> str:
    return "Аудит" if normalize_report_lang(lang) == "ru" else "Audit"


def report_docx_labels(lang: str) -> Dict[str, str | List[str]]:
    if normalize_report_lang(lang) == "ru":
        return {
            "status": "Статус",
            "issues": "Ошибок",
            "headers": ["ID ошибки", "Человеческая локация", "Проблема", "Рекомендация"],
        }
    return {
        "status": "Status",
        "issues": "Issues",
        "headers": ["Issue ID", "Human Location", "Problem", "Recommendation"],
    }


REPORT_VALUE_MAP_RU = {
    "paragraph": "абзац",
    "table_cell": "ячейка таблицы",
    "slide_text": "текст слайда",
    "safe_replace": "безопасная замена",
    "requires_review": "требуется проверка",
    "none": "без правки",
    "high": "высокая",
    "medium": "средняя",
    "low": "низкая",
    "review": "проверка",
    "ok": "успешно",
}


def localized_report_value(value: Any, report_lang: str) -> Any:
    if normalize_report_lang(report_lang) != "ru":
        return value
    text = str(value or "")
    return REPORT_VALUE_MAP_RU.get(text.strip().lower(), value)


def localized_human_text(value: Any, report_lang: str) -> str:
    text = str(value or "")
    if normalize_report_lang(report_lang) == "ru":
        text = re.sub(r"\bCHECK\s*:", "ПРОВЕРКА:", text, flags=re.IGNORECASE)
    return text


def formatted_technical_location(value: Any, report_lang: str) -> str:
    technical = value if isinstance(value, dict) else {}
    if normalize_report_lang(report_lang) != "ru":
        return json.dumps(technical, ensure_ascii=False)
    labels = {
        "part": "часть",
        "xpath_hint": "XPath",
        "paragraph_index": "абзац",
        "table_index": "таблица",
        "row_index": "строка",
        "cell_index": "ячейка",
        "slide_index": "слайд",
    }
    return "; ".join(
        f"{labels.get(key, key)}: {item}" for key, item in technical.items() if item not in (None, "")
    )


CYRILLIC_WORD_RE = re.compile(r"[А-Яа-яЁё]{2,}")
LATIN_WORD_RE = re.compile(r"[A-Za-z]{3,}")
QUOTED_SPAN_RE = re.compile(r"«[^»]*»|„[^“”]*[“”]|“[^”]*”|\"[^\"]*\"|'[^']*'|`[^`]*`")
LATIN_TOKEN_RE = re.compile(r"[A-Za-z][A-Za-z0-9_.\-]*")
TECHNICAL_LATIN_RE = re.compile(
    r"[A-Za-z]*[a-z][A-Z][A-Za-z]*"
    r"|[A-Z]{2,}"
    r"|[A-Za-z]+[0-9_][A-Za-z0-9_]*"
    r"|[A-Za-z]+(?:\.[A-Za-z0-9]+)+"
)
LANGUAGE_REPAIR_VERSION = 2


def strip_untranslatable_latin(text: str, quote: str = "") -> str:
    """Убрать цитаты и технические токены: они остаются латиницей и не означают английский текст."""
    cleaned = QUOTED_SPAN_RE.sub(" ", text)
    quoted_tokens = {token for token in LATIN_TOKEN_RE.findall(str(quote or "")) if len(token) >= 3}
    for token in sorted(quoted_tokens, key=len, reverse=True):
        cleaned = cleaned.replace(token, " ")
    return TECHNICAL_LATIN_RE.sub(" ", cleaned)


def human_text_needs_russian_repair(text: Any, quote: Any = "") -> bool:
    value = str(text or "").strip()
    if not value:
        return False
    if re.search(r"\bCHECK\s*:", value, flags=re.IGNORECASE):
        return True
    prose = strip_untranslatable_latin(value, str(quote or ""))
    latin_words = LATIN_WORD_RE.findall(prose)
    cyrillic_words = CYRILLIC_WORD_RE.findall(value)
    if not cyrillic_words:
        return len(latin_words) >= 2
    return len(latin_words) >= 3 and len(latin_words) > len(cyrillic_words)


def invalid_human_report_fields(issues: List[Dict[str, Any]], report_lang: str) -> List[Dict[str, Any]]:
    if normalize_report_lang(report_lang) != "ru":
        return []
    invalid: List[Dict[str, Any]] = []
    for issue_index, issue in enumerate(issues):
        quote = str(issue.get("quote") or "")
        for field in ("problem", "recommendation"):
            text = str(issue.get(field) or "").strip()
            if human_text_needs_russian_repair(text, quote):
                invalid.append(
                    {
                        "target_id": f"L{len(invalid) + 1:03d}",
                        "issue_index": issue_index,
                        "issue_id": str(issue.get("issue_id") or "?"),
                        "field": field,
                        "text": text,
                        "quote": quote,
                    }
                )
    return invalid


def validate_human_report_language(issues: List[Dict[str, Any]], report_lang: str) -> None:
    """Reject English-only or English-dominant human explanations in a Russian report."""
    invalid = invalid_human_report_fields(issues, report_lang)
    if invalid:
        sample = ", ".join(f"{item['issue_id']}.{item['field']}" for item in invalid[:8])
        raise RuntimeError(
            "Russian report contains English-only or English-dominant human-readable fields: "
            f"{sample}. The LLM output language requirement was not satisfied."
        )


def sha256_file(path: Path) -> str:
    h = hashlib.sha256()
    with path.open("rb") as fh:
        for chunk in iter(lambda: fh.read(1024 * 1024), b""):
            h.update(chunk)
    return h.hexdigest()


def write_json(path: Path, payload: Dict[str, Any]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")


def read_json(path: Path) -> Dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8-sig", errors="replace"))


def artifact_paths(paths, source: Path) -> Dict[str, Path]:
    parent = rel_parent(source, paths.input_dir)
    stem = source.stem
    return {
        "pdf": paths.work_dir / "rendered_pdf" / parent / f"{stem}.pdf",
        "marked": paths.work_dir / "marked_ooxml" / parent / f"{stem}__marked{source.suffix.lower()}",
        "pages": paths.work_dir / "extracted_pdf_text" / parent / f"{stem}__pages.json",
        "render_log": paths.logs_dir / parent / f"{stem}__render.json",
        "block_map": paths.logs_dir / parent / f"{stem}__block_map.json",
        "render_map": paths.logs_dir / parent / f"{stem}__render_map.json",
        "audit": paths.logs_dir / parent / f"{stem}__audit.json",
        "language_repair": paths.logs_dir / parent / f"{stem}__language_repair.json",
        "annotation_log": paths.logs_dir / parent / f"{stem}__annotation.json",
        "fix_log": paths.logs_dir / parent / f"{stem}__fix.json",
        "issues_in": paths.cache_dir / parent / f"{stem}__issues.json",
        "table": paths.output_dir / parent / f"{stem}__audit_table.xlsx",
        "report": paths.output_dir / parent / f"{stem}__audit_report.docx",
        "annotated": paths.output_dir / parent / f"{stem}__annotated{source.suffix.lower()}",
        "fixed": paths.output_dir / parent / f"{stem}__fixed{source.suffix.lower()}",
    }


def cmd_scan(paths, recursive: bool) -> int:
    ensure_project_dirs(paths)
    docs = iter_documents(paths.input_dir, recursive=recursive)
    payload = {
        "input_dir": str(paths.input_dir),
        "recursive": bool(recursive),
        "supported_extensions": [".docx", ".pptx"],
        "ignored_temp_prefix": "~$",
        "documents": [
            {
                "relative_path": source_relative(doc, paths.input_dir).as_posix(),
                "extension": doc.suffix.lower(),
                "size": doc.stat().st_size,
            }
            for doc in docs
        ],
    }
    write_json(paths.logs_dir / "scan.json", payload)
    print(f"[SCAN] input={paths.input_dir}")
    print(f"[SCAN] documents={len(docs)}")
    for item in payload["documents"]:
        print(f"  - {item['relative_path']}")
    return 0


def render_one(paths, source: Path, renderer: str) -> Dict[str, Any]:
    if renderer != "com":
        raise RuntimeError(f"Unsupported renderer: {renderer}")

    outputs = artifact_paths(paths, source)
    started = time.perf_counter()
    block_map = build_block_map(source, paths.input_dir)
    write_json(outputs["block_map"], block_map)
    create_marked_copy(source, outputs["marked"], block_map)

    render_payload: Dict[str, Any] = {
        "source_relative_path": source_relative(source, paths.input_dir).as_posix(),
        "source_sha256": sha256_file(source),
        "renderer": "microsoft_office_com",
        "office_app": "Word.Application" if source.suffix.lower() == ".docx" else "PowerPoint.Application",
        "rendered_pdf": str(outputs["pdf"]),
        "marked_copy": str(outputs["marked"]),
        "status": "started",
        "duration_seconds": None,
        "warnings": [],
    }

    try:
        warnings = export_office_document_to_pdf(paths.root, outputs["marked"], outputs["pdf"], outputs["render_log"].with_suffix(".log"))
        render_payload["warnings"] = warnings
        pages_payload = write_pdf_pages(outputs["pdf"], outputs["pages"])
        render_map = build_render_map(block_map, pages_payload)
        write_json(outputs["render_map"], render_map)
        render_payload["status"] = "ok"
        render_payload["pages_json"] = str(outputs["pages"])
        render_payload["render_map"] = str(outputs["render_map"])
    except Exception as exc:
        render_payload["status"] = "failed"
        render_payload["error"] = str(exc)
        fallback = build_render_map(block_map, {"pages": []})
        write_json(outputs["render_map"], fallback)
    finally:
        render_payload["duration_seconds"] = round(time.perf_counter() - started, 3)
        write_json(outputs["render_log"], render_payload)

    return render_payload


def cmd_render(paths, recursive: bool, renderer: str) -> int:
    ensure_project_dirs(paths)
    docs = iter_documents(paths.input_dir, recursive=recursive)
    if not docs:
        print(f"[INFO] No .docx/.pptx files found in: {paths.input_dir}")
        return 0

    failures = 0
    for index, source in enumerate(docs, start=1):
        label = source_relative(source, paths.input_dir).as_posix()
        started = time.perf_counter()
        print(f"[RENDER] [{index}/{len(docs)}] {label}")
        payload = render_one(paths, source, renderer)
        elapsed = format_duration(time.perf_counter() - started)
        if payload.get("status") == "ok":
            print(f"[OK] {label} ({elapsed})")
        else:
            failures += 1
            print(f"[FAIL] {label}: {payload.get('error', 'unknown error')} ({elapsed})")
    return 2 if failures else 0


def load_reusable_render(outputs: Dict[str, Path], source: Path) -> Dict[str, Any] | None:
    try:
        if not outputs["render_log"].exists() or not outputs["block_map"].exists() or not outputs["render_map"].exists():
            return None
        payload = read_json(outputs["render_log"])
        if payload.get("status") != "ok":
            return None
        if payload.get("source_sha256") != sha256_file(source):
            return None
        return payload
    except Exception:
        return None


def normalize_issues(
    raw_issues: List[Dict[str, Any]],
    block_map: Dict[str, Any],
    render_map: Dict[str, Any],
) -> List[Dict[str, Any]]:
    blocks = {block["block_id"]: block for block in block_map.get("blocks", []) or []}
    renders = render_entry_by_block(render_map)
    normalized: List[Dict[str, Any]] = []

    for index, raw in enumerate(raw_issues, start=1):
        block_id = str(raw.get("block_id") or "").strip()
        if block_id not in blocks:
            continue
        block = blocks[block_id]
        if should_suppress_false_positive_issue(raw, block):
            continue
        render_entry = renders.get(block_id, {})
        technical = block.get("technical_location", {}) or {}
        issue_id = str(raw.get("issue_id") or f"E{index:03d}")
        normalized.append(
            {
                "issue_id": issue_id,
                "severity": raw.get("severity", "review"),
                "rule_id": raw.get("rule_id", ""),
                "human_location": build_human_location(block, render_entry),
                "page": render_entry.get("page"),
                "object_type": block.get("object_type") or block.get("kind", ""),
                "table": technical.get("table_index"),
                "row": technical.get("row_index"),
                "cell": technical.get("cell_index"),
                "paragraph": technical.get("paragraph_index"),
                "quote": raw.get("quote", block.get("text", "")),
                "problem": raw.get("problem") or raw.get("violation", ""),
                "recommendation": raw.get("recommendation") or raw.get("fix", ""),
                "fix_mode": raw.get("fix_mode", "requires_review"),
                "old_text": raw.get("old_text") or raw.get("source_text") or raw.get("replace_old", ""),
                "new_text": raw.get("new_text") or raw.get("replacement_text") or raw.get("replace_new", ""),
                "confidence": raw.get("confidence", ""),
                "status": raw.get("status", "requires_review"),
                "block_id": block_id,
                "technical_location": technical,
                "render_location": render_entry,
                "page_source": render_entry.get("page_source"),
                "page_confidence": render_entry.get("page_confidence"),
                "llm_raw_location": raw.get("location", ""),
            }
        )
    return normalized


def should_suppress_false_positive_issue(raw: Dict[str, Any], block: Dict[str, Any]) -> bool:
    return is_official_nested_quote_false_positive(raw, str(block.get("text") or ""))


def is_official_nested_quote_false_positive(raw: Dict[str, Any], block_text: str) -> bool:
    """Filter LLM-only quote complaints against exact official titles/citations."""
    problem = str(raw.get("problem") or raw.get("violation") or "")
    recommendation = str(raw.get("recommendation") or raw.get("fix") or "")
    combined = f"{problem} {recommendation}".lower()
    if "кавыч" not in combined:
        return False
    if not any(marker in combined for marker in ("сломан", "закрыва", "открыва", "незакрыт", "пар")):
        return False

    text = " ".join(str(block_text or "").split())
    if "«" not in text or "»" not in text:
        return False
    if not OFFICIAL_TITLE_MARKERS_RE.search(text):
        return False

    if text.count("«") == text.count("»"):
        return True

    return bool(OFFICIAL_NESTED_SOURCE_RE.search(text))


def build_empty_audit(source: Path, block_map: Dict[str, Any], render_map: Dict[str, Any], status: str) -> Dict[str, Any]:
    return {
        "source_relative_path": block_map.get("source_relative_path", source.name),
        "status": status,
        "issues": [],
        "meta": {
            "issue_source": "none",
            "note": "No LLM issue list was supplied to the pipeline; report artifacts are generated for smoke/structure validation.",
        },
    }


DEFAULT_AUDIT_REASONING = {
    "openai": "high",
    "gemini": "medium",
    "xai": "high",
    "anthropic": "medium",
}


def configured_active_provider() -> str:
    try:
        settings = load_settings()
        provider = str(settings.get("active_provider") or "openai").strip().lower()
        return provider if provider in {"openai", "gemini", "xai", "anthropic", "sidecar"} else "openai"
    except Exception:
        return "openai"


def resolve_audit_model(provider: str, model: str = "", settings: Dict[str, Any] | None = None) -> str:
    selected = str(model or "").strip()
    if selected:
        return selected
    provider = str(provider or "").strip().lower()
    settings = settings or load_settings()
    if provider == "openai":
        return resolve_model("openai", "audit", settings)
    if provider == "gemini":
        return resolve_model("gemini", "audit_fast", settings)
    if provider == "xai":
        return resolve_model("xai", "audit", settings)
    if provider == "anthropic":
        return resolve_model("anthropic", "audit", settings)
    return ""


def load_raw_issues(path: Path) -> tuple[List[Dict[str, Any]], str]:
    if not path.exists():
        return [], "none"
    payload = read_json(path)
    if isinstance(payload.get("issues"), list):
        return payload["issues"], str(path)
    if isinstance(payload.get("rows"), list):
        return payload["rows"], str(path)
    raise RuntimeError(f"Invalid issue sidecar, expected issues[] or rows[]: {path}")


def block_text_for_llm(block: Dict[str, Any]) -> str:
    text = str(block.get("text") or "").strip()
    block_id = str(block.get("block_id") or block.get("id") or "").strip()
    kind = str(block.get("object_type") or block.get("kind") or "").strip()
    technical = block.get("technical_location", {}) or {}
    paragraph = technical.get("paragraph_index")
    table = technical.get("table_index")
    row = technical.get("row_index")
    cell = technical.get("cell_index")
    loc_bits = []
    if kind:
        loc_bits.append(f"type={kind}")
    if paragraph is not None:
        loc_bits.append(f"paragraph={paragraph}")
    if table is not None:
        loc_bits.append(f"table={table}")
    if row is not None:
        loc_bits.append(f"row={row}")
    if cell is not None:
        loc_bits.append(f"cell={cell}")
    suffix = f" ({', '.join(loc_bits)})" if loc_bits else ""
    return f"[BLOCK:{block_id}]{suffix}\n{text}"


def blocks_for_llm(block_map: Dict[str, Any]) -> List[Dict[str, Any]]:
    out: List[Dict[str, Any]] = []
    for block in block_map.get("blocks", []) or []:
        block_id = str(block.get("block_id") or "").strip()
        text = str(block.get("text") or "").strip()
        if not block_id or not text:
            continue
        out.append(
            {
                "id": block_id,
                "block_id": block_id,
                "kind": block.get("object_type") or block.get("kind", ""),
                "location": f"[BLOCK:{block_id}]",
                "text": block_text_for_llm(block),
            }
        )
    return out


def row_with_block_id(row: Dict[str, Any], block_map: Dict[str, Any]) -> Dict[str, Any]:
    import re

    out = dict(row)
    if out.get("block_id"):
        return out

    location = str(out.get("location") or "")
    match = re.search(r"\[BLOCK:([^\]]+)\]", location)
    if match:
        out["block_id"] = match.group(1).strip()
        return out

    quote = str(out.get("quote") or "").strip()
    if quote:
        for block in block_map.get("blocks", []) or []:
            text = str(block.get("text") or "")
            if quote in text:
                out["block_id"] = block.get("block_id")
                return out
    return out


def llm_cache_path(outputs: Dict[str, Path], provider: str, model: str, reasoning: str) -> Path:
    safe_model = "".join(ch if ch.isalnum() or ch in {"-", "_"} else "_" for ch in model)[:80] or "default"
    safe_reasoning = "".join(ch if ch.isalnum() or ch in {"-", "_"} else "_" for ch in reasoning)[:32] or "default"
    base = outputs["issues_in"]
    return base.with_name(f"{base.stem}__{provider}_{safe_model}_{safe_reasoning}{base.suffix}")


def sha256_text(text: str) -> str:
    return hashlib.sha256(str(text or "").encode("utf-8", errors="ignore")).hexdigest()


def llm_cache_signature(
    *,
    source_sha256: str,
    provider: str,
    model: str,
    reasoning: str,
    rules_context: str,
    instructions: str,
    chunk_tokens: int,
    overlap_tokens: int,
    min_chunks: int,
    max_output_tokens: int,
) -> Dict[str, Any]:
    return {
        "cache_version": 2,
        "source_sha256": str(source_sha256 or ""),
        "provider": str(provider or "").lower().strip(),
        "model": str(model or "").strip(),
        "reasoning": str(reasoning or "").strip(),
        "rules_sha256": sha256_text(rules_context),
        "instructions_sha256": sha256_text(instructions),
        "chunk_tokens": int(chunk_tokens),
        "overlap_tokens": int(overlap_tokens),
        "min_chunks": int(min_chunks),
        "max_output_tokens": int(max_output_tokens),
    }


def llm_cache_mismatch_reason(payload: Dict[str, Any], expected: Dict[str, Any]) -> str:
    actual = payload.get("cache_signature")
    if not isinstance(actual, dict):
        return "missing cache_signature"
    for key, expected_value in expected.items():
        if actual.get(key) != expected_value:
            return f"{key} changed"
    return ""


LLM_USAGE_TOKEN_KEYS = ("input_tokens", "output_tokens", "total_tokens", "reasoning_tokens")


def new_llm_usage_total() -> Dict[str, Any]:
    return {"calls": 0, "input_tokens": 0, "output_tokens": 0, "total_tokens": 0, "reasoning_tokens": 0}


def add_llm_usage(usage_total: Dict[str, Any], usage: Dict[str, Any], service_tier: str = "") -> None:
    usage_total["calls"] = int(usage_total.get("calls", 0) or 0) + 1
    for key in LLM_USAGE_TOKEN_KEYS:
        usage_total[key] = int(usage_total.get(key, 0) or 0) + int(usage.get(key, 0) or 0)
    if service_tier:
        usage_total["service_tier"] = service_tier


def merge_llm_usage_total(usage_total: Dict[str, Any], aggregate: Dict[str, Any]) -> None:
    usage_total["calls"] = int(usage_total.get("calls", 0) or 0) + int(aggregate.get("calls", 0) or 0)
    for key in LLM_USAGE_TOKEN_KEYS:
        usage_total[key] = int(usage_total.get(key, 0) or 0) + int(aggregate.get(key, 0) or 0)
    if aggregate.get("service_tier"):
        usage_total["service_tier"] = aggregate["service_tier"]


def _chunk_cache_sort_key(key: str) -> int:
    try:
        return int(key)
    except Exception:
        return 0


def rows_from_chunk_cache(chunks: Dict[str, Any]) -> List[Dict[str, Any]]:
    rows: List[Dict[str, Any]] = []
    for key in sorted(chunks.keys(), key=_chunk_cache_sort_key):
        item = chunks.get(key)
        if not isinstance(item, dict):
            continue
        for row in item.get("rows") or []:
            if isinstance(row, dict):
                rows.append(row)
    return rows


def usage_from_chunk_cache(chunks: Dict[str, Any]) -> Dict[str, Any]:
    usage_total = new_llm_usage_total()
    for key in sorted(chunks.keys(), key=_chunk_cache_sort_key):
        item = chunks.get(key)
        if not isinstance(item, dict):
            continue
        usage = item.get("usage") if isinstance(item.get("usage"), dict) else {}
        add_llm_usage(usage_total, usage, str(item.get("service_tier") or ""))
    return usage_total


def llm_cache_payload(
    *,
    provider: str,
    model: str,
    reasoning: str,
    doc_hash: str,
    cache_signature: Dict[str, Any],
    source_relative_path: str,
    rows: List[Dict[str, Any]],
    usage_total: Dict[str, Any],
    status: str,
    total_chunks: int,
    chunk_cache: Dict[str, Any],
) -> Dict[str, Any]:
    return {
        "provider": provider,
        "model": model,
        "reasoning": reasoning,
        "source_sha256": doc_hash,
        "cache_signature": cache_signature,
        "source_relative_path": source_relative_path,
        "status": status,
        "total_chunks": int(total_chunks),
        "completed_chunks": len(chunk_cache),
        "rows": rows,
        "usage": usage_total,
        "chunks": chunk_cache,
    }


def write_llm_progress_cache(
    cache_path: Path,
    *,
    provider: str,
    model: str,
    reasoning: str,
    doc_hash: str,
    cache_signature: Dict[str, Any],
    source_relative_path: str,
    rows: List[Dict[str, Any]],
    usage_total: Dict[str, Any],
    total_chunks: int,
    chunk_cache: Dict[str, Any],
) -> None:
    write_json(
        cache_path,
        llm_cache_payload(
            provider=provider,
            model=model,
            reasoning=reasoning,
            doc_hash=doc_hash,
            cache_signature=cache_signature,
            source_relative_path=source_relative_path,
            rows=dedupe_rows(rows),
            usage_total=usage_total,
            status="partial",
            total_chunks=total_chunks,
            chunk_cache=chunk_cache,
        ),
    )


def usage_summary(usage: Dict[str, Any]) -> str:
    total = int(usage.get("total_tokens", 0) or 0)
    reasoning = int(usage.get("reasoning_tokens", 0) or 0)
    input_tokens = int(usage.get("input_tokens", 0) or 0)
    output_tokens = int(usage.get("output_tokens", 0) or 0)
    return f"tokens={total} (reasoning={reasoning}) in={input_tokens} out={output_tokens}"


def chunk_eta(durations: List[float], total_chunks: int, completed: int, workers: int = 1) -> str:
    if not durations or completed >= total_chunks:
        return "00:00"
    avg = sum(durations) / len(durations)
    lanes = max(1, int(workers or 1))
    return format_duration(avg * (total_chunks - completed) / lanes)


def _language_repair_signature(
    targets: List[Dict[str, Any]], provider: str, model: str, report_lang: str
) -> str:
    payload = {
        "version": LANGUAGE_REPAIR_VERSION,
        "provider": provider,
        "model": model,
        "report_lang": normalize_report_lang(report_lang),
        "targets": [
            {
                "target_id": item["target_id"],
                "issue_id": item["issue_id"],
                "field": item["field"],
                "text": item["text"],
            }
            for item in targets
        ],
    }
    raw = json.dumps(payload, ensure_ascii=False, sort_keys=True, separators=(",", ":"))
    return hashlib.sha256(raw.encode("utf-8")).hexdigest()


def _apply_language_repairs(
    issues: List[Dict[str, Any]],
    targets: List[Dict[str, Any]],
    replacements: List[Dict[str, Any]],
) -> List[Dict[str, Any]]:
    target_by_id = {item["target_id"]: item for item in targets}
    replacement_by_id: Dict[str, str] = {}
    for item in replacements:
        target_id = str(item.get("target_id") or "").strip()
        text = str(item.get("text") or item.get("repaired_text") or "").strip()
        if target_id in replacement_by_id:
            raise RuntimeError(f"Language repair returned duplicate target_id: {target_id}")
        if target_id in target_by_id and text:
            replacement_by_id[target_id] = text

    missing = [target_id for target_id in target_by_id if target_id not in replacement_by_id]
    if missing:
        raise RuntimeError(f"Language repair omitted target(s): {', '.join(missing)}")

    repaired = [dict(issue) for issue in issues]
    for target_id, target in target_by_id.items():
        repaired[target["issue_index"]][target["field"]] = replacement_by_id[target_id]
    return repaired


def _call_language_repair_provider(
    *,
    provider: str,
    model: str,
    instructions: str,
    user_prompt: str,
    repair_hash: str,
    max_output_tokens: int,
    timeout_sec: float,
    max_retries: int,
    service_tier: str,
) -> tuple[Dict[str, Any], Dict[str, Any], str, str]:
    settings = load_settings()
    provider = provider.strip().lower()
    if provider == "openai":
        from openai import OpenAI
        from providers.openai_provider import call_json_object

        api_key = resolve_api_key("openai", settings)
        if not api_key:
            raise RuntimeError("OpenAI API key not found for language repair.")
        client = OpenAI(api_key=api_key, timeout=timeout_sec, max_retries=0)
        obj, usage, tier = call_json_object(
            client,
            model=model,
            instructions=instructions,
            user_prompt=user_prompt,
            reasoning_effort="low",
            max_output_tokens=max_output_tokens,
            timeout_sec=timeout_sec,
            max_retries=max_retries,
            service_tier=service_tier,
            use_idempotency=True,
            doc_hash=repair_hash,
            chunk_index=0,
        )
        return obj, usage, tier, "low"
    if provider == "gemini":
        from google import genai
        from providers.gemini_provider import call_structured

        api_key = resolve_api_key("gemini", settings)
        if not api_key:
            raise RuntimeError("Gemini API key not found for language repair.")
        client = genai.Client(api_key=api_key)
        obj, usage, tier = call_structured(
            client,
            model=model,
            system_instruction=instructions,
            user_prompt=user_prompt,
            thinking_level="minimal",
            max_retries=max_retries,
        )
        return obj, usage, tier, "minimal"
    if provider == "xai":
        from providers.xai_provider import call_json_object

        api_key = resolve_api_key("xai", settings)
        if not api_key:
            raise RuntimeError("xAI API key not found for language repair.")
        obj, usage, tier = call_json_object(
            api_key=api_key,
            model=model,
            instructions=instructions,
            user_prompt=user_prompt,
            max_output_tokens=max_output_tokens,
            timeout_sec=timeout_sec,
            max_retries=max_retries,
            use_idempotency=True,
            doc_hash=repair_hash,
            chunk_index=0,
            reasoning_effort="low",
        )
        return obj, usage, tier, "low"
    if provider == "anthropic":
        from providers.anthropic_provider import call_json_object

        api_key = resolve_api_key("anthropic", settings)
        if not api_key:
            raise RuntimeError("Anthropic API key not found for language repair.")
        obj, usage, tier = call_json_object(
            api_key=api_key,
            model=model,
            instructions=instructions,
            user_prompt=user_prompt,
            max_output_tokens=max_output_tokens,
            timeout_sec=timeout_sec,
            max_retries=max_retries,
            reasoning_effort="low",
        )
        return obj, usage, tier, "low"
    raise RuntimeError(f"Language repair is not available for provider: {provider}")


def language_repair_provider_chain(provider: str, model: str) -> List[Dict[str, str]]:
    provider = str(provider or "").strip().lower()
    chain = [{"provider": provider, "model": str(model or "").strip()}]
    fallback_provider = {
        "openai": "gemini",
        "gemini": "openai",
        "xai": "gemini",
        "anthropic": "openai",
    }.get(provider, "")
    if fallback_provider:
        fallback_model = resolve_audit_model(fallback_provider)
        if fallback_model:
            chain.append({"provider": fallback_provider, "model": fallback_model})
    return chain


def repair_human_report_language(
    issues: List[Dict[str, Any]],
    *,
    report_lang: str,
    provider: str,
    model: str,
    output_path: Path,
    max_output_tokens: int,
    timeout_sec: float,
    max_retries: int,
    service_tier: str,
) -> tuple[List[Dict[str, Any]], Dict[str, Any]]:
    targets = invalid_human_report_fields(issues, report_lang)
    if not targets:
        return issues, {"status": "not_needed", "fields": 0, "cache_hit": False}
    if provider.strip().lower() == "sidecar":
        validate_human_report_language(issues, report_lang)

    public_targets = [
        {
            "target_id": item["target_id"],
            "issue_id": item["issue_id"],
            "field": item["field"],
            "text": item["text"],
            "quote": item["quote"],
        }
        for item in targets
    ]
    signature = _language_repair_signature(targets, provider, model, report_lang)
    if output_path.exists():
        cached = read_json(output_path)
        cached_replacements = cached.get("replacements") if isinstance(cached.get("replacements"), list) else []
        if cached.get("status") == "repaired" and cached.get("signature") == signature and cached_replacements:
            try:
                repaired = _apply_language_repairs(issues, targets, cached_replacements)
                validate_human_report_language(repaired, report_lang)
                print(f"[LANGUAGE REPAIR] cache hit fields={len(targets)}")
                return repaired, {
                    "status": "repaired",
                    "fields": len(targets),
                    "cache_hit": True,
                    "usage": {},
                    "reasoning": str(cached.get("repair_reasoning") or cached.get("reasoning") or ""),
                    "provider": str(cached.get("repair_provider") or cached.get("provider") or provider),
                    "model": str(cached.get("repair_model") or cached.get("model") or model),
                    "fallback_used": bool(cached.get("fallback_used")),
                }
            except Exception:
                pass

    instructions = (
        "Ты исправляешь только язык человекочитаемых полей русского отчёта. "
        "Перепиши каждый переданный text на естественном русском языке, сохранив точный смысл. "
        "Не меняй target_id, issue_id, факты, числа, имена собственные и точные цитаты. "
        "Не добавляй пояснений и не используй префикс CHECK. "
        "Верни только JSON: {\"rows\":[{\"target_id\":\"L001\",\"text\":\"русский текст\"}]}"
    )
    user_prompt = json.dumps(
        {"request": "Return JSON only with repaired Russian text fields.", "targets": public_targets},
        ensure_ascii=False,
        indent=2,
    )
    repair_max_output = max(800, min(int(max_output_tokens), max(1200, len(targets) * 350)))
    pending = {
        "status": "pending",
        "version": LANGUAGE_REPAIR_VERSION,
        "signature": signature,
        "provider": provider,
        "model": model,
        "targets": public_targets,
    }
    write_json(output_path, pending)
    attempts: List[Dict[str, Any]] = []
    usage_total = new_llm_usage_total()
    chain = language_repair_provider_chain(provider, model)
    last_error: Exception | None = None
    for attempt_index, repair_target in enumerate(chain, start=1):
        repair_provider = repair_target["provider"]
        repair_model = repair_target["model"]
        attempt_hash = hashlib.sha256(
            f"{signature}|{repair_provider}|{repair_model}".encode("utf-8")
        ).hexdigest()
        started = time.perf_counter()
        response: Dict[str, Any] = {}
        repair_usage: Dict[str, Any] = {}
        tier = ""
        repair_reasoning = ""
        print(
            f"[LANGUAGE REPAIR] attempt={attempt_index}/{len(chain)} "
            f"provider={repair_provider} model={repair_model} fields={len(targets)} start"
        )
        try:
            response, repair_usage, tier, repair_reasoning = _call_language_repair_provider(
                provider=repair_provider,
                model=repair_model,
                instructions=instructions,
                user_prompt=user_prompt,
                repair_hash=attempt_hash,
                max_output_tokens=repair_max_output,
                timeout_sec=timeout_sec,
                max_retries=max_retries,
                service_tier=service_tier,
            )
            add_llm_usage(usage_total, repair_usage, tier)
            replacements = response.get("rows") or response.get("repairs") or []
            if not isinstance(replacements, list):
                replacements = []
            replacements = [item for item in replacements if isinstance(item, dict)]
            repaired = _apply_language_repairs(issues, targets, replacements)
            validate_human_report_language(repaired, report_lang)
            elapsed = time.perf_counter() - started
            attempts.append(
                {
                    "attempt": attempt_index,
                    "status": "repaired",
                    "provider": repair_provider,
                    "model": repair_model,
                    "reasoning": repair_reasoning,
                    "usage": repair_usage,
                    "elapsed_sec": round(elapsed, 3),
                    "replacements": replacements,
                }
            )
            diagnostic = {
                **pending,
                "status": "repaired",
                "repair_provider": repair_provider,
                "repair_model": repair_model,
                "repair_reasoning": repair_reasoning,
                "fallback_used": attempt_index > 1,
                "replacements": replacements,
                "usage": usage_total,
                "attempts": attempts,
            }
            write_json(output_path, diagnostic)
            print(
                f"[LANGUAGE REPAIR] ok attempt={attempt_index}/{len(chain)} "
                f"time={format_duration(elapsed)} {usage_summary(repair_usage)} fields={len(targets)}"
            )
            return repaired, {
                "status": "repaired",
                "fields": len(targets),
                "cache_hit": False,
                "usage": usage_total,
                "service_tier": tier,
                "reasoning": repair_reasoning,
                "provider": repair_provider,
                "model": repair_model,
                "fallback_used": attempt_index > 1,
            }
        except Exception as exc:
            last_error = exc
            elapsed = time.perf_counter() - started
            attempts.append(
                {
                    "attempt": attempt_index,
                    "status": "failed",
                    "provider": repair_provider,
                    "model": repair_model,
                    "reasoning": repair_reasoning,
                    "usage": repair_usage,
                    "elapsed_sec": round(elapsed, 3),
                    "error": str(exc),
                    "response": response,
                }
            )
            failed = dict(pending)
            failed.update(
                {
                    "status": "retrying" if attempt_index < len(chain) else "failed",
                    "error": str(exc),
                    "usage": usage_total,
                    "attempts": attempts,
                }
            )
            write_json(output_path, failed)
            if attempt_index < len(chain):
                next_target = chain[attempt_index]
                print(
                    f"[LANGUAGE REPAIR] attempt={attempt_index} failed: {exc}; "
                    f"fallback={next_target['provider']} model={next_target['model']}"
                )

    raise RuntimeError(
        f"Language repair failed for {len(targets)} field(s) after {len(chain)} provider attempt(s); "
        f"no audit report was published: {last_error}"
    ) from last_error


def run_chunk_jobs(
    pairs: List[Any],
    *,
    label: str,
    workers: int,
    chunk_cache: Dict[str, Any],
    call: Any,
    extract_rows: Any,
    persist: Any = None,
) -> None:
    """Fill chunk_cache for every chunk that is not cached yet, `workers` at a time."""
    total = len(pairs)
    pending: List[tuple[int, str, str]] = []
    for index, (overlap, chunk) in enumerate(pairs, start=1):
        cached = chunk_cache.get(str(index))
        if cached:
            print(f"  [CHUNK {index}/{total}] {label} cache hit rows={len(cached.get('rows') or [])}")
            continue
        pending.append((index, overlap, chunk))

    if not pending:
        return

    lanes = min(resolve_workers(workers), len(pending))
    if lanes > 1:
        print(f"[LLM] {label} concurrency={lanes} chunks at a time")

    lock = threading.Lock()
    durations: List[float] = []

    def work(job: tuple[int, str, str]) -> None:
        index, overlap, chunk = job
        print(f"  [CHUNK {index}/{total}] {label} start")
        started = time.perf_counter()
        obj, usage, tier = call(index, overlap, chunk)
        elapsed = time.perf_counter() - started
        chunk_rows = extract_rows(obj)
        with lock:
            durations.append(elapsed)
            chunk_cache[str(index)] = {
                "index": index,
                "rows": chunk_rows,
                "usage": usage,
                "service_tier": tier,
                "elapsed_sec": round(elapsed, 3),
            }
            if persist is not None:
                persist(chunk_cache)
            eta = chunk_eta(durations, total, len(chunk_cache), lanes)
        print(
            f"  [CHUNK {index}/{total}] {label} ok "
            f"time={format_duration(elapsed)} {usage_summary(usage)} "
            f"eta={eta} rows={len(chunk_rows)}"
        )

    run_jobs(pending, work, workers=lanes)


def collect_llm_issues(
    source: Path,
    outputs: Dict[str, Path],
    block_map: Dict[str, Any],
    provider: str,
    model: str,
    reasoning: str,
    chunk_tokens: int,
    overlap_tokens: int,
    min_chunks: int,
    max_retries: int,
    max_output_tokens: int,
    timeout_sec: float,
    service_tier: str,
    resume: bool,
    report_lang: str = "en",
    workers: int = DEFAULT_WORKERS,
) -> tuple[List[Dict[str, Any]], str, Dict[str, Any]]:
    provider = provider.lower().strip()
    settings = load_settings()
    model = resolve_audit_model(provider, model, settings)
    if not model:
        raise RuntimeError(
            f"No model selected for provider: {provider}. Choose a model in the GUI dropdown, "
            "or add/check a favorite model so config\\gui_model_cache.json can provide the fallback."
        )

    cache_path = llm_cache_path(outputs, provider, model, reasoning)
    doc_hash = sha256_file(source)
    rules_context = load_rules()
    instructions = build_instructions(rules_context, report_lang=report_lang)
    cache_signature = llm_cache_signature(
        source_sha256=doc_hash,
        provider=provider,
        model=model,
        reasoning=reasoning,
        rules_context=rules_context,
        instructions=instructions,
        chunk_tokens=chunk_tokens,
        overlap_tokens=overlap_tokens,
        min_chunks=min_chunks,
        max_output_tokens=max_output_tokens,
    )
    source_relative_path = str(block_map.get("source_relative_path", source.name))
    partial_payload: Dict[str, Any] = {}
    if resume and cache_path.exists():
        payload = read_json(cache_path)
        stale_reason = llm_cache_mismatch_reason(payload, cache_signature)
        if not stale_reason:
            status = str(payload.get("status") or "complete").strip().lower()
            if status == "partial":
                partial_payload = payload
                chunks = payload.get("chunks") if isinstance(payload.get("chunks"), dict) else {}
                print(f"[LLM] provider={provider} model={model} cache=partial chunks={len(chunks)}; resuming")
            else:
                rows = payload.get("rows") or payload.get("issues") or []
                usage = payload.get("usage", {}) or {}
                print(f"[LLM] provider={provider} model={model} cache=hit rows={len(rows)} {usage_summary(usage)}")
                return rows, str(cache_path), payload.get("usage", {})
        if stale_reason:
            print(f"[LLM] provider={provider} model={model} cache=stale reason={stale_reason}; recomputing")

    blocks = blocks_for_llm(block_map)
    pairs = build_chunks(
        blocks,
        chunk_tokens=chunk_tokens,
        overlap_tokens=overlap_tokens,
        min_chunks=min_chunks,
        max_input_tokens=0,
        prompt_overhead_tokens=3000,
    )
    print(f"[LLM] provider={provider} model={model} chunks={len(pairs)}")

    rows: List[Dict[str, Any]] = []
    usage_total = new_llm_usage_total()
    durations: List[float] = []
    chunk_cache: Dict[str, Any] = {}

    if partial_payload:
        cached_chunks = partial_payload.get("chunks") if isinstance(partial_payload.get("chunks"), dict) else {}
        for key, item in cached_chunks.items():
            index = _chunk_cache_sort_key(str(key))
            if 1 <= index <= len(pairs) and isinstance(item, dict):
                chunk_cache[str(index)] = item
        if chunk_cache:
            rows = rows_from_chunk_cache(chunk_cache)
            usage_total = usage_from_chunk_cache(chunk_cache)
            print(
                f"[LLM] provider={provider} model={model} resume "
                f"chunks={len(chunk_cache)}/{len(pairs)} rows={len(rows)} {usage_summary(usage_total)}"
            )

    def extract_rows(obj: Dict[str, Any]) -> List[Dict[str, Any]]:
        chunk_rows = obj.get("rows") or obj.get("issues") or []
        if not isinstance(chunk_rows, list):
            chunk_rows = []
        return [row_with_block_id(row, block_map) for row in chunk_rows if isinstance(row, dict)]

    def persist(snapshot: Dict[str, Any]) -> None:
        write_llm_progress_cache(
            cache_path,
            provider=provider,
            model=model,
            reasoning=reasoning,
            doc_hash=doc_hash,
            cache_signature=cache_signature,
            source_relative_path=source_relative_path,
            rows=rows_from_chunk_cache(snapshot),
            usage_total=usage_from_chunk_cache(snapshot),
            total_chunks=len(pairs),
            chunk_cache=snapshot,
        )

    if provider == "openai":
        from openai import OpenAI
        from providers.openai_provider import call_json_object

        api_key = resolve_api_key("openai", settings)
        if not api_key:
            raise RuntimeError("OpenAI API key not found. Set OPENAI_API_KEY or config\\api_key_openai.txt.")
        client = OpenAI(api_key=api_key, timeout=timeout_sec, max_retries=0)
        label = "OpenAI"

        def call(index: int, overlap: str, chunk: str):
            return call_json_object(
                client,
                model=model,
                instructions=instructions,
                user_prompt=build_user_prompt(overlap, chunk),
                reasoning_effort=reasoning,
                max_output_tokens=max_output_tokens,
                timeout_sec=timeout_sec,
                max_retries=max_retries,
                service_tier=service_tier,
                use_idempotency=True,
                doc_hash=doc_hash,
                chunk_index=index,
            )

    elif provider == "gemini":
        from google import genai
        from providers.gemini_provider import call_structured

        api_key = resolve_api_key("gemini", settings)
        if not api_key:
            raise RuntimeError("Gemini API key not found. Set GEMINI_API_KEY or config\\api_key_gemini.txt.")
        client = genai.Client(api_key=api_key)
        label = "Gemini"

        def call(index: int, overlap: str, chunk: str):
            return call_structured(
                client,
                model=model,
                system_instruction=instructions,
                user_prompt=build_user_prompt(overlap, chunk),
                thinking_level=reasoning,
                max_retries=max_retries,
            )

    elif provider == "xai":
        from providers.xai_provider import call_json_object

        api_key = resolve_api_key("xai", settings)
        if not api_key:
            raise RuntimeError("xAI API key not found. Set XAI_API_KEY or config\\api_key_xai.txt.")
        label = "xAI"

        def call(index: int, overlap: str, chunk: str):
            return call_json_object(
                api_key=api_key,
                model=model,
                instructions=instructions,
                user_prompt=build_user_prompt(overlap, chunk),
                max_output_tokens=max_output_tokens,
                timeout_sec=timeout_sec,
                max_retries=max_retries,
                use_idempotency=True,
                doc_hash=doc_hash,
                chunk_index=index,
                reasoning_effort=reasoning,
            )

    elif provider == "anthropic":
        from providers.anthropic_provider import call_json_object

        api_key = resolve_api_key("anthropic", settings)
        if not api_key:
            raise RuntimeError("Anthropic API key not found. Set ANTHROPIC_API_KEY or config\\api_key_anthropic.txt.")
        label = "Claude"

        def call(index: int, overlap: str, chunk: str):
            return call_json_object(
                api_key=api_key,
                model=model,
                instructions=instructions,
                user_prompt=build_user_prompt(overlap, chunk),
                max_output_tokens=max_output_tokens,
                timeout_sec=timeout_sec,
                max_retries=max_retries,
                reasoning_effort=reasoning,
            )

    else:
        raise RuntimeError(f"Unsupported LLM provider: {provider}")

    run_chunk_jobs(
        pairs,
        label=label,
        workers=workers,
        chunk_cache=chunk_cache,
        call=call,
        extract_rows=extract_rows,
        persist=persist if resume else None,
    )

    rows = dedupe_rows(rows_from_chunk_cache(chunk_cache))
    usage_total = usage_from_chunk_cache(chunk_cache)
    payload = llm_cache_payload(
        provider=provider,
        model=model,
        reasoning=reasoning,
        doc_hash=doc_hash,
        cache_signature=cache_signature,
        source_relative_path=source_relative_path,
        rows=rows,
        usage_total=usage_total,
        status="complete",
        total_chunks=len(pairs),
        chunk_cache=chunk_cache,
    )
    write_json(cache_path, payload)
    return rows, str(cache_path), usage_total


def strict_render_map_failed(render_map: Dict[str, Any]) -> bool:
    entries = render_map.get("entries", []) or []
    return bool(entries) and not any(entry.get("page") for entry in entries)


def audit_one(
    paths,
    source: Path,
    renderer: str,
    require_render_map: bool,
    apply_fixes: str,
    report_lang: str,
    provider: str,
    model: str,
    reasoning: str,
    chunk_tokens: int,
    overlap_tokens: int,
    min_chunks: int,
    max_retries: int,
    max_output_tokens: int,
    timeout_sec: float,
    service_tier: str,
    resume: bool,
    workers: int = DEFAULT_WORKERS,
) -> Dict[str, Any]:
    outputs = artifact_paths(paths, source)
    render_payload = load_reusable_render(outputs, source)
    if render_payload:
        print("[RENDER] reuse existing COM render-map")
    else:
        render_payload = render_one(paths, source, renderer)
    block_map = read_json(outputs["block_map"])
    render_map = read_json(outputs["render_map"])

    if require_render_map and (render_payload.get("status") != "ok" or strict_render_map_failed(render_map)):
        audit_payload = build_empty_audit(source, block_map, render_map, "failed_render_map")
        write_json(outputs["audit"], audit_payload)
        return audit_payload

    if provider == "sidecar":
        raw_issues, issue_source = load_raw_issues(outputs["issues_in"])
        usage = {}
    else:
        raw_issues, issue_source, usage = collect_llm_issues(
            source,
            outputs,
            block_map,
            provider,
            model,
            reasoning,
            chunk_tokens,
            overlap_tokens,
            min_chunks,
            max_retries,
            max_output_tokens,
            timeout_sec,
            service_tier,
            resume,
            report_lang,
            workers,
        )
    issues = normalize_issues(raw_issues, block_map, render_map)
    issues, language_repair = repair_human_report_language(
        issues,
        report_lang=report_lang,
        provider=provider,
        model=model,
        output_path=outputs["language_repair"],
        max_output_tokens=max_output_tokens,
        timeout_sec=timeout_sec,
        max_retries=max_retries,
        service_tier=service_tier,
    )
    repair_usage = language_repair.get("usage") if isinstance(language_repair.get("usage"), dict) else {}
    if repair_usage and not language_repair.get("cache_hit"):
        merge_llm_usage_total(usage, repair_usage)
    validate_human_report_language(issues, report_lang)
    audit_payload = {
        "source_relative_path": block_map.get("source_relative_path", source.name),
        "status": "ok",
        "issues": issues,
        "meta": {
            "issue_source": issue_source,
            "provider": provider,
            "model": model,
            "reasoning": reasoning,
            "usage": usage,
            "language_repair": {
                "status": language_repair.get("status"),
                "fields": int(language_repair.get("fields", 0) or 0),
                "cache_hit": bool(language_repair.get("cache_hit")),
                "reasoning": language_repair.get("reasoning", ""),
                "provider": language_repair.get("provider", ""),
                "model": language_repair.get("model", ""),
                "fallback_used": bool(language_repair.get("fallback_used")),
                "diagnostic": str(outputs["language_repair"]) if language_repair.get("fields") else "",
            },
            "note": "Final locations are normalized by Python from block_id, block_map, and render_map.",
        },
    }
    if provider == "xai":
        from providers.xai_provider import assert_clean_json_text

        assert_clean_json_text(audit_payload)
    write_json(outputs["audit"], audit_payload)
    write_audit_table(outputs["table"], audit_payload["issues"], report_lang=report_lang)
    write_audit_docx(outputs["report"], source.name, audit_payload, report_lang=report_lang)
    write_annotated_document(source, outputs["annotated"], block_map, audit_payload["issues"], outputs["annotation_log"])

    if apply_fixes == "safe":
        shutil.copy2(source, outputs["fixed"])
        write_json(outputs["fix_log"], {"status": "no_safe_fixes_available", "fixed_document": str(outputs["fixed"])})

    return audit_payload


def cmd_audit(
    paths,
    recursive: bool,
    renderer: str,
    require_render_map: bool,
    apply_fixes: str,
    report_lang: str,
    provider: str,
    model: str,
    reasoning: str,
    chunk_tokens: int,
    overlap_tokens: int,
    min_chunks: int,
    max_retries: int,
    max_output_tokens: int,
    timeout_sec: float,
    service_tier: str,
    resume: bool,
    workers: int = DEFAULT_WORKERS,
) -> int:
    ensure_project_dirs(paths)
    model = resolve_audit_model(provider, model)
    workers = resolve_workers(workers)
    print(
        f"[AUDIT CONFIG] provider={provider} model={model or '<none>'} "
        f"reasoning={reasoning} report_lang={normalize_report_lang(report_lang)} workers={workers}"
    )
    docs = iter_documents(paths.input_dir, recursive=recursive)
    if not docs:
        print(f"[INFO] No .docx/.pptx files found in: {paths.input_dir}")
        return 0

    failures = 0
    for index, source in enumerate(docs, start=1):
        rel = source_relative(source, paths.input_dir).as_posix()
        print(f"[AUDIT] [{index}/{len(docs)}] {rel}")
        payload = audit_one(
            paths,
            source,
            renderer,
            require_render_map,
            apply_fixes,
            report_lang,
            provider,
            model,
            reasoning,
            chunk_tokens,
            overlap_tokens,
            min_chunks,
            max_retries,
            max_output_tokens,
            timeout_sec,
            service_tier,
            resume,
            workers,
        )
        if payload.get("status") == "failed_render_map":
            failures += 1
            print(f"[FAIL] render-map required but unavailable: {rel}")
        else:
            print(f"[OK] issues={len(payload.get('issues', []))}")
    return 2 if failures else 0


def write_audit_table(out_path: Path, issues: List[Dict[str, Any]], report_lang: str = "en") -> None:
    from openpyxl import Workbook
    from openpyxl.styles import Alignment, Border, Font, PatternFill, Side

    out_path.parent.mkdir(parents=True, exist_ok=True)
    wb = Workbook()
    ws = wb.active
    ws.title = report_sheet_title(report_lang)
    ws.append(report_columns(report_lang))

    for issue in issues:
        technical = issue.get("technical_location", {})
        ws.append(
            [
                issue.get("issue_id", ""),
                issue.get("human_location", ""),
                issue.get("page", ""),
                localized_report_value(issue.get("object_type", ""), report_lang),
                issue.get("table", ""),
                issue.get("row", ""),
                issue.get("cell", ""),
                issue.get("paragraph", ""),
                issue.get("quote", ""),
                localized_human_text(issue.get("problem", ""), report_lang),
                localized_human_text(issue.get("recommendation", ""), report_lang),
                localized_report_value(issue.get("fix_mode", ""), report_lang),
                localized_report_value(issue.get("confidence", ""), report_lang),
                issue.get("old_text", ""),
                issue.get("new_text", ""),
                issue.get("block_id", ""),
                formatted_technical_location(technical, report_lang),
            ]
        )

    header_fill = PatternFill("solid", fgColor="1F4E78")
    header_font = Font(name="Tahoma", size=9, bold=True, color="FFFFFF")
    text_font = Font(name="Tahoma", size=9, color="1F1F1F")
    even_fill = PatternFill("solid", fgColor="F7F3EA")
    odd_fill = PatternFill("solid", fgColor="EFE7D8")
    thin = Side(style="thin", color="D9E2F3")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)

    for row in ws.iter_rows():
        for cell in row:
            cell.border = border
            cell.alignment = Alignment(horizontal="left", vertical="center", wrap_text=True, indent=1)
            if cell.row == 1:
                cell.font = header_font
                cell.fill = header_fill
                cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
            else:
                cell.font = text_font
                cell.fill = even_fill if cell.row % 2 == 0 else odd_fill

    for column_cells in ws.columns:
        max_len = max(len(str(cell.value or "")) for cell in column_cells)
        letter = column_cells[0].column_letter
        ws.column_dimensions[letter].width = min(max(max_len + 2, 8), 62)

    ws.freeze_panes = "A2"
    ws.auto_filter.ref = ws.dimensions
    ws.sheet_view.showGridLines = False

    for letter in ["I", "J", "K", "M"]:
        ws.column_dimensions[letter].width = max(ws.column_dimensions[letter].width or 0, 42)

    wb.save(str(out_path))


def write_audit_docx(out_path: Path, title: str, audit_payload: Dict[str, Any], report_lang: str = "en") -> None:
    from docx import Document
    from docx.shared import Pt

    out_path.parent.mkdir(parents=True, exist_ok=True)
    issues = audit_payload.get("issues", []) or []
    doc = Document()
    style = doc.styles["Normal"]
    style.font.name = "Tahoma"
    style.font.size = Pt(10)
    labels = report_docx_labels(report_lang)
    doc.add_heading(title, 0)
    doc.add_paragraph(
        f"{labels['status']}: {localized_report_value(audit_payload.get('status', ''), report_lang)}"
    )
    doc.add_paragraph(f"{labels['issues']}: {len(issues)}")

    table = doc.add_table(rows=1, cols=4)
    table.style = "Table Grid"
    headers = labels["headers"]
    for idx, header in enumerate(headers):
        table.cell(0, idx).text = header
    for issue in issues:
        row = table.add_row().cells
        row[0].text = str(issue.get("issue_id", ""))
        row[1].text = str(issue.get("human_location", ""))
        row[2].text = localized_human_text(issue.get("problem", ""), report_lang)
        row[3].text = localized_human_text(issue.get("recommendation", ""), report_lang)
    doc.save(str(out_path))


def write_annotated_document(
    source: Path,
    out_path: Path,
    block_map: Dict[str, Any],
    issues: List[Dict[str, Any]],
    log_path: Path,
) -> None:
    out_path.parent.mkdir(parents=True, exist_ok=True)
    issue_by_block: Dict[str, List[Dict[str, Any]]] = {}
    for issue in issues:
        issue_by_block.setdefault(issue.get("block_id", ""), []).append(issue)

    if source.suffix.lower() == ".docx":
        refs_payload = block_paragraph_refs_docx(source)
        doc = refs_payload["document"]
        refs = refs_payload["refs"]
        for block_id, block_issues in issue_by_block.items():
            paragraph = refs.get(block_id)
            if paragraph is None:
                continue
            for issue in block_issues:
                run = paragraph.add_run(f" ⟦{issue.get('issue_id', '')}⟧")
                run.bold = True
                try:
                    from docx.shared import RGBColor

                    run.font.color.rgb = RGBColor(0xC0, 0x56, 0x00)
                except Exception:
                    pass
        doc.save(str(out_path))
    elif source.suffix.lower() == ".pptx":
        refs_payload = block_paragraph_refs_pptx(source)
        prs = refs_payload["presentation"]
        refs = refs_payload["refs"]
        for block_id, block_issues in issue_by_block.items():
            paragraph = refs.get(block_id)
            if paragraph is None:
                continue
            for issue in block_issues:
                run = paragraph.add_run()
                run.text = f" ⟦{issue.get('issue_id', '')}⟧"
                try:
                    from pptx.dml.color import RGBColor

                    run.font.color.rgb = RGBColor(0xC0, 0x56, 0x00)
                    run.font.bold = True
                except Exception:
                    pass
        prs.save(str(out_path))
    else:
        shutil.copy2(source, out_path)

    write_json(
        log_path,
        {
            "source": str(source),
            "annotated": str(out_path),
            "issues": len(issues),
            "issue_ids": [issue.get("issue_id", "") for issue in issues],
            "status": "ok",
        },
    )


def _strip_anchor_runs(runs: Iterable[Any]) -> int:
    removed = 0
    for run in runs:
        text = str(getattr(run, "text", "") or "")
        matches = ERROR_ANCHOR_RE.findall(text)
        if not matches:
            continue
        run.text = ERROR_ANCHOR_RE.sub("", text)
        removed += len(matches)
    return removed


def _iter_docx_paragraphs(container: Any) -> Iterable[Any]:
    for paragraph in getattr(container, "paragraphs", []) or []:
        yield paragraph
    for table in getattr(container, "tables", []) or []:
        for row in table.rows:
            for cell in row.cells:
                yield from _iter_docx_paragraphs(cell)


def _iter_pptx_paragraphs(shapes: Any) -> Iterable[Any]:
    for shape in shapes:
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        yield paragraph
            continue
        if getattr(shape, "has_text_frame", False):
            for paragraph in shape.text_frame.paragraphs:
                yield paragraph
        nested_shapes = getattr(shape, "shapes", None)
        if nested_shapes is not None:
            yield from _iter_pptx_paragraphs(nested_shapes)


def strip_error_anchors_from_document(source: Path, out_path: Path) -> int:
    out_path.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(source, out_path)
    suffix = source.suffix.lower()
    removed = 0

    if suffix == ".docx":
        from docx import Document

        doc = Document(str(out_path))
        for paragraph in _iter_docx_paragraphs(doc):
            removed += _strip_anchor_runs(paragraph.runs)
        doc.save(str(out_path))
        return removed

    if suffix == ".pptx":
        from pptx import Presentation

        prs = Presentation(str(out_path))
        for slide in prs.slides:
            for paragraph in _iter_pptx_paragraphs(slide.shapes):
                removed += _strip_anchor_runs(paragraph.runs)
        prs.save(str(out_path))
        return removed

    return removed


def unanchored_output_path(source: Path) -> Path:
    stem = source.stem
    if stem.endswith("__annotated"):
        stem = stem[: -len("__annotated")]
    return source.with_name(f"{stem}__unanchored{source.suffix.lower()}")


def cmd_strip_anchors(paths, recursive: bool, all_documents: bool) -> int:
    root = paths.output_dir
    pattern = "**/*" if recursive else "*"
    docs = [
        path
        for path in sorted(root.glob(pattern))
        if path.is_file()
        and path.suffix.lower() in {".docx", ".pptx"}
        and not path.name.startswith("~$")
        and not path.stem.endswith("__unanchored")
        and (all_documents or path.stem.endswith("__annotated"))
    ]
    if not docs:
        print(f"[INFO] No annotated DOCX/PPTX files found in: {root}")
        return 0

    total = 0
    for source in docs:
        out_path = unanchored_output_path(source)
        removed = strip_error_anchors_from_document(source, out_path)
        total += removed
        print(f"[ANCHORS] {source.relative_to(root).as_posix()} -> {out_path.relative_to(root).as_posix()} removed={removed}")
    print(f"[DONE] unanchored={len(docs)} anchors_removed={total}")
    return 0


def iter_audit_logs(logs_root: Path) -> Iterable[Path]:
    return sorted(logs_root.rglob("*__audit.json"))


def source_from_audit(paths, audit_payload: Dict[str, Any]) -> Path:
    if paths.input_dir.is_file():
        return paths.input_dir
    return paths.input_dir / audit_payload["source_relative_path"]


def cmd_report(paths, from_logs: Path, report_lang: str) -> int:
    count = 0
    for audit_log in iter_audit_logs(from_logs):
        payload = read_json(audit_log)
        source = source_from_audit(paths, payload)
        outputs = artifact_paths(paths, source)
        write_audit_table(outputs["table"], payload.get("issues", []) or [], report_lang=report_lang)
        write_audit_docx(outputs["report"], source.name, payload, report_lang=report_lang)
        count += 1
        print(f"[REPORT] {outputs['table']}")
        print(f"[REPORT] {outputs['report']}")
    print(f"[DONE] reports={count}")
    return 0


def cmd_annotate(paths, from_logs: Path) -> int:
    count = 0
    for audit_log in iter_audit_logs(from_logs):
        payload = read_json(audit_log)
        source = source_from_audit(paths, payload)
        outputs = artifact_paths(paths, source)
        block_map = read_json(outputs["block_map"])
        write_annotated_document(source, outputs["annotated"], block_map, payload.get("issues", []) or [], outputs["annotation_log"])
        count += 1
        print(f"[ANNOTATE] {outputs['annotated']}")
    print(f"[DONE] annotated={count}")
    return 0


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Audion Docs AI recursive DOCX/PPTX pipeline.")
    sub = parser.add_subparsers(dest="command", required=True)

    p_scan = sub.add_parser("scan")
    p_scan.add_argument("--recursive", action="store_true", default=True)
    p_scan.add_argument("--no-recursive", dest="recursive", action="store_false")

    p_render = sub.add_parser("render")
    p_render.add_argument("--recursive", action="store_true", default=True)
    p_render.add_argument("--no-recursive", dest="recursive", action="store_false")
    p_render.add_argument("--renderer", choices=["com"], default="com")

    p_audit = sub.add_parser("audit")
    p_audit.add_argument("--recursive", action="store_true", default=True)
    p_audit.add_argument("--no-recursive", dest="recursive", action="store_false")
    p_audit.add_argument("--renderer", choices=["com"], default="com")
    p_audit.add_argument("--require-render-map", action="store_true")
    p_audit.add_argument("--apply-fixes", choices=["none", "safe"], default="none")
    p_audit.add_argument("--report-lang", choices=["en", "ru"], default=normalize_report_lang())
    p_audit.add_argument(
        "--provider",
        choices=["openai", "gemini", "xai", "anthropic", "sidecar"],
        default=configured_active_provider(),
    )
    p_audit.add_argument("--model", default="")
    p_audit.add_argument("--reasoning", default="")
    p_audit.add_argument("--chunk-tokens", type=int, default=12000)
    p_audit.add_argument("--overlap-tokens", type=int, default=1200)
    p_audit.add_argument("--min-chunks", type=int, default=4)
    p_audit.add_argument("--max-retries", type=int, default=3)
    p_audit.add_argument("--max-output-tokens", type=int, default=32000)
    p_audit.add_argument("--timeout-sec", type=float, default=420.0)
    p_audit.add_argument("--service-tier", default="default")
    p_audit.add_argument("--resume", action="store_true")
    p_audit.add_argument(
        "--workers",
        type=int,
        default=DEFAULT_WORKERS,
        help=f"chunks processed at a time (1-{MAX_WORKERS})",
    )

    p_annotate = sub.add_parser("annotate")
    p_annotate.add_argument("--from-logs", default="logs")

    p_strip = sub.add_parser("strip-anchors")
    p_strip.add_argument("--recursive", action="store_true", default=True)
    p_strip.add_argument("--no-recursive", dest="recursive", action="store_false")
    p_strip.add_argument("--all-documents", action="store_true")

    p_report = sub.add_parser("report")
    p_report.add_argument("--from-logs", default="logs")
    p_report.add_argument("--report-lang", choices=["en", "ru"], default=normalize_report_lang())

    return parser


def main() -> int:
    # The pipeline is often launched with stdout redirected into the NiceGUI
    # terminal or an external supervisor.  Force line-by-line delivery even
    # when the caller forgot ``python -u``/PYTHONUNBUFFERED, otherwise chunk
    # progress remains invisible until the whole audit finishes.
    for stream in (sys.stdout, sys.stderr):
        reconfigure = getattr(stream, "reconfigure", None)
        if callable(reconfigure):
            reconfigure(line_buffering=True, write_through=True)

    root = Path(__file__).resolve().parents[1]
    paths = default_project_paths(root)
    parser = build_parser()
    args = parser.parse_args()

    if args.command == "scan":
        return cmd_scan(paths, args.recursive)
    if args.command == "render":
        return cmd_render(paths, args.recursive, args.renderer)
    if args.command == "audit":
        return cmd_audit(
            paths,
            args.recursive,
            args.renderer,
            args.require_render_map,
            args.apply_fixes,
            args.report_lang,
            args.provider,
            args.model,
            args.reasoning or DEFAULT_AUDIT_REASONING.get(args.provider, "selected"),
            args.chunk_tokens,
            args.overlap_tokens,
            args.min_chunks,
            args.max_retries,
            args.max_output_tokens,
            args.timeout_sec,
            args.service_tier,
            args.resume,
            args.workers,
        )
    if args.command == "annotate":
        return cmd_annotate(paths, root / args.from_logs)
    if args.command == "strip-anchors":
        return cmd_strip_anchors(paths, args.recursive, args.all_documents)
    if args.command == "report":
        return cmd_report(paths, root / args.from_logs, args.report_lang)
    raise RuntimeError(f"Unknown command: {args.command}")


if __name__ == "__main__":
    raise SystemExit(main())
